from __future__ import annotations

import re
from collections import Counter
from io import BytesIO
from pathlib import Path
from typing import Callable, Dict


_PAGE_PATTERNS = [
    re.compile(r"^\s*(sayfa|page)\s*\d+\s*(/|of)?\s*\d*\s*$", re.IGNORECASE),
    re.compile(r"^\s*\d+\s*/\s*\d+\s*$"),
    re.compile(r"^\s*\d+\s*of\s*\d+\s*$", re.IGNORECASE),
]
_ONLY_NUMBER_LINE = re.compile(r"^\s*\d{1,3}\s*$")


def _is_page_artifact_line(line: str) -> bool:
    s = (line or "").strip()
    if not s:
        return False
    return any(p.match(s) for p in _PAGE_PATTERNS)


def clean_extracted_text(text: str, min_repeated: int = 3) -> str:
    if not text:
        return ""

    lines = [ln.rstrip() for ln in text.splitlines()]
    kept = [ln for ln in lines if not _is_page_artifact_line(ln)]

    normalized = [re.sub(r"\s+", " ", ln.strip().lower()) for ln in kept]
    counts = Counter(s for s in normalized if 0 < len(s) <= 60)

    out: list[str] = []
    for original, norm in zip(kept, normalized):
        if 0 < len(norm) <= 60 and counts[norm] >= min_repeated:
            continue
        if original.strip():
            out.append(original.strip())
    return "\n".join(out)


def _load_pdf(path: str) -> str:
    import fitz

    parts: list[str] = []
    try:
        doc = fitz.open(path)
        try:
            for page in doc:
                t = page.get_text("text")
                if t:
                    parts.append(t)
        finally:
            doc.close()
    except Exception as e:
        raise RuntimeError(f"PDF okunamadı: {e}") from e

    return clean_extracted_text("\n".join(parts).strip())


def _load_pptx(path: str) -> str:
    from pptx import Presentation
    from pptx.enum.shapes import PP_PLACEHOLDER

    skip = {PP_PLACEHOLDER.SLIDE_NUMBER, PP_PLACEHOLDER.DATE, PP_PLACEHOLDER.FOOTER}
    try:
        prs = Presentation(path)
    except Exception as e:
        raise RuntimeError(f"PPTX okunamadı: {e}") from e

    out: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.is_placeholder:
                try:
                    if shape.placeholder_format.type in skip:
                        continue
                except Exception:
                    pass
            if hasattr(shape, "text") and shape.text:
                out.append(shape.text)

    return clean_extracted_text("\n".join(out).strip())


def _load_docx(path: str) -> str:
    from docx import Document

    with open(path, "rb") as f:
        doc = Document(BytesIO(f.read()))

    parts = [(p.text or "").strip() for p in doc.paragraphs]
    return "\n".join(p for p in parts if p)


def _load_txt(path: str) -> str:
    try:
        with open(path, "r", encoding="utf-8") as f:
            return f.read()
    except UnicodeDecodeError:
        for enc in ("latin5", "iso-8859-9", "windows-1254"):
            try:
                with open(path, "r", encoding=enc) as f:
                    return f.read()
            except Exception:
                continue
    raise RuntimeError("TXT dosyası okunamadı (encoding sorunu)")


_LOADERS: Dict[str, Callable[[str], str]] = {
    ".pdf": _load_pdf,
    ".pptx": _load_pptx,
    ".docx": _load_docx,
    ".txt": _load_txt,
}


def load_file(path: str) -> str:
    ext = Path(path).suffix.lower()
    loader = _LOADERS.get(ext)
    if loader is None:
        raise ValueError(f"Desteklenmeyen dosya türü: {ext}")
    return loader(path)
